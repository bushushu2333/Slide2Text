from pptx import Presentation
from PIL import Image
import pytesseract
import os

def extract_ppt_content(ppt_file, output_dir="output"):
    """
    从PPT文件中提取内容并保存为结构化数据
    :param ppt_file: PPT文件路径
    :param output_dir: 输出目录
    :return: 结构化数据
    """
    # 创建输出目录
    os.makedirs(output_dir, exist_ok=True)
    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    ppt_pages = []
    presentation = Presentation(ppt_file)

    for i, slide in enumerate(presentation.slides):
        page_data = {
            "page_index": i + 1,
            "raw_text": "",
            "ocr_text": "",
            "images": [],
            "metadata": {
                "slide_title": "",
                "notes": "",
                "slide_layout": "",
                "slide_number": i + 1,
                "author": ""
            }
        }

        # 提取文本内容
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                page_data["raw_text"] += shape.text + "\n"

        # 提取图片并进行OCR识别
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "image"):
                # 保存图片
                image_path = os.path.join(images_dir, f"page{i+1}_image{j+1}.png")
                with open(image_path, "wb") as f:
                    f.write(shape.image.blob)
                
                # OCR识别
                ocr_result = perform_ocr(image_path)
                
                # 更新页面OCR文本
                if ocr_result:
                    page_data["ocr_text"] += ocr_result + "\n"

                # 添加图片信息
                page_data["images"].append({
                    "image_id": j + 1,
                    "image_path": image_path,
                    "ocr_result": ocr_result,
                    "description": f"Page {i+1} Image {j+1}"
                })

        # 提取元数据
        if slide.has_notes_slide:
            page_data["metadata"]["notes"] = slide.notes_slide.notes_text_frame.text
        if slide.shapes.title:
            page_data["metadata"]["slide_title"] = slide.shapes.title.text

        ppt_pages.append(page_data)

    return ppt_pages

def perform_ocr(image_path):
    """
    对图片进行OCR识别
    :param image_path: 图片路径
    :return: 识别结果文本
    """
    try:
        # 打开图片并进行预处理
        image = Image.open(image_path)
        # 转换为灰度图
        image = image.convert('L')
        # 使用pytesseract进行OCR识别，指定中文语言
        text = pytesseract.image_to_string(image, lang='chi_sim')
        return text.strip()
    except Exception as e:
        print(f"OCR识别失败: {str(e)}")
        return ""

def save_structured_data(data, output_file):
    """
    保存结构化数据到文件
    :param data: 结构化数据
    :param output_file: 输出文件路径
    """
    import json
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

